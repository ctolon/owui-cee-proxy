#!/usr/bin/env python3
"""
Generate the binary fixtures used by the integration test suite.

Run with:
    python3 -m pip install --user openpyxl python-docx python-pptx \
        ebooklib pillow olefile
    python3 scripts/gen-fixtures.py

The output goes to test/integration/fixtures/. Each file is small and
self-describing where possible (subject lines / titles say
"owui-cee-proxy routing fixture").

Formats produced:
    sample.txt   sample.html   sample.md       sample.csv
    sample.json  sample.xml    sample.rtf      sample.eml
    sample.pdf   sample.png    sample.jpeg
    sample.docx  sample.xlsx   sample.pptx
    sample.odt   sample.epub   sample.msg

The PDF, MSG and PNG are hand-crafted byte by byte so the fixture set
stays self-contained (no system tooling like libreoffice required).
"""
from __future__ import annotations

import io
import json
import os
import struct
import sys
import zipfile
from email.message import EmailMessage
from pathlib import Path
from xml.sax.saxutils import escape

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "test" / "integration" / "fixtures"


def _write(name: str, data: bytes) -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    path = OUT / name
    path.write_bytes(data)
    print(f"  wrote {path.relative_to(ROOT)} ({len(data)} bytes)")


# ── plain-text-ish formats ────────────────────────────────────────────


def gen_txt() -> None:
    _write("sample.txt", b"hello world\n\nowui-cee-proxy routing fixture (txt).\n")


def gen_html() -> None:
    body = (
        "<!doctype html><html><head><meta charset=utf-8>"
        "<title>owui-cee-proxy routing fixture</title></head>"
        "<body><h1>HTML fixture</h1>"
        "<p>This document exercises the proxy's <code>text/html</code>"
        " routing rule.</p></body></html>\n"
    )
    _write("sample.html", body.encode())


def gen_md() -> None:
    _write(
        "sample.md",
        b"# owui-cee-proxy routing fixture\n\n"
        b"Markdown body for `text/markdown` routing tests.\n",
    )


def gen_csv() -> None:
    _write(
        "sample.csv",
        b"id,name,score\n1,alice,7.5\n2,bob,9.1\n3,carol,8.4\n",
    )


def gen_json() -> None:
    payload = {"fixture": "owui-cee-proxy", "purpose": "routing", "n": 3}
    _write("sample.json", json.dumps(payload, indent=2).encode())


def gen_xml() -> None:
    body = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        "<doc><meta><title>owui-cee-proxy fixture</title></meta>"
        "<body><p>XML body.</p></body></doc>\n"
    )
    _write("sample.xml", body.encode())


def gen_rtf() -> None:
    rtf = (
        r"{\rtf1\ansi\deff0 "
        r"{\fonttbl{\f0\fnil\fcharset0 Helvetica;}}"
        r"\f0\fs24 owui-cee-proxy routing fixture (RTF). \par "
        r"This file exercises the proxy's application/rtf routing.}"
    )
    _write("sample.rtf", rtf.encode("utf-8"))


def gen_eml() -> None:
    msg = EmailMessage()
    msg["Subject"] = "owui-cee-proxy routing fixture"
    msg["From"] = "alice@example.com"
    msg["To"] = "bob@example.com"
    msg.set_content(
        "This is a real RFC 822 email body.\n"
        "It exercises message/rfc822 routing.\n"
    )
    _write("sample.eml", bytes(msg))


# ── PDF (hand-crafted minimal) ────────────────────────────────────────


def gen_pdf() -> None:
    # 1-page Helvetica "Hello, owui-cee-proxy!" PDF.
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
        b"/Resources<</Font<</F1 5 0 R>>>>/Contents 4 0 R>>endobj\n"
        b"4 0 obj<</Length 53>>stream\n"
        b"BT /F1 18 Tf 100 720 Td (Hello, owui-cee-proxy!) Tj ET\n"
        b"endstream endobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f\n"
        b"0000000010 00000 n\n"
        b"0000000054 00000 n\n"
        b"0000000101 00000 n\n"
        b"0000000194 00000 n\n"
        b"0000000291 00000 n\n"
        b"trailer<</Size 6/Root 1 0 R>>\nstartxref\n352\n%%EOF\n"
    )
    _write("sample.pdf", pdf)


# ── images ────────────────────────────────────────────────────────────


def gen_png_and_jpeg() -> None:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (320, 80), "white")
    draw = ImageDraw.Draw(img)
    draw.text((10, 30), "owui-cee-proxy routing fixture", fill="black")

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    _write("sample.png", buf.getvalue())

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85)
    _write("sample.jpeg", buf.getvalue())


# ── Office Open XML (real DOCX/XLSX/PPTX) ─────────────────────────────


def gen_docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("owui-cee-proxy routing fixture", level=1)
    doc.add_paragraph(
        "This DOCX exercises the proxy's "
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document "
        "routing rule."
    )
    doc.add_paragraph("Engine: Docling.")
    buf = io.BytesIO()
    doc.save(buf)
    _write("sample.docx", buf.getvalue())


def gen_xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "fixture"
    ws.append(["id", "name", "score"])
    ws.append([1, "alice", 7.5])
    ws.append([2, "bob", 9.1])
    ws.append([3, "carol", 8.4])
    buf = io.BytesIO()
    wb.save(buf)
    _write("sample.xlsx", buf.getvalue())


def gen_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "owui-cee-proxy routing fixture"
    slide.placeholders[1].text = (
        "PPTX exercises the proxy's "
        "application/vnd.openxmlformats-officedocument.presentationml.presentation routing."
    )
    buf = io.BytesIO()
    prs.save(buf)
    _write("sample.pptx", buf.getvalue())


# ── ODT (hand-crafted ZIP — avoids LibreOffice dependency) ────────────


def gen_odt() -> None:
    content_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<office:document-content '
        'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
        'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
        'xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0" '
        'office:version="1.2">'
        '<office:body><office:text>'
        '<text:h text:outline-level="1">owui-cee-proxy routing fixture</text:h>'
        '<text:p>ODT body — exercises application/vnd.oasis.opendocument.text routing.</text:p>'
        '</office:text></office:body></office:document-content>'
    )
    manifest_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">'
        '<manifest:file-entry manifest:full-path="/" '
        'manifest:media-type="application/vnd.oasis.opendocument.text"/>'
        '<manifest:file-entry manifest:full-path="content.xml" '
        'manifest:media-type="text/xml"/>'
        '</manifest:manifest>'
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        # mimetype must be the FIRST entry, stored uncompressed.
        zi = zipfile.ZipInfo("mimetype")
        zi.compress_type = zipfile.ZIP_STORED
        z.writestr(zi, "application/vnd.oasis.opendocument.text")
        z.writestr("META-INF/manifest.xml", manifest_xml)
        z.writestr("content.xml", content_xml)
    _write("sample.odt", buf.getvalue())


# ── EPUB (hand-crafted ZIP) ───────────────────────────────────────────


def gen_epub() -> None:
    content_opf = """<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" unique-identifier="bookid" version="3.0">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:identifier id="bookid">urn:uuid:owui-cee-proxy-fixture</dc:identifier>
    <dc:title>owui-cee-proxy routing fixture</dc:title>
    <dc:language>en</dc:language>
    <meta property="dcterms:modified">2026-05-10T00:00:00Z</meta>
  </metadata>
  <manifest>
    <item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/>
    <item id="c1" href="ch1.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine>
    <itemref idref="c1"/>
  </spine>
</package>
"""
    nav_xhtml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<html xmlns="http://www.w3.org/1999/xhtml" xmlns:epub="http://www.idpf.org/2007/ops">'
        "<head><title>nav</title></head>"
        '<body><nav epub:type="toc"><ol><li><a href="ch1.xhtml">Chapter 1</a></li></ol></nav></body></html>'
    )
    ch1_xhtml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<html xmlns="http://www.w3.org/1999/xhtml"><head><title>Chapter 1</title></head>'
        "<body><h1>owui-cee-proxy routing fixture</h1>"
        "<p>EPUB body — exercises application/epub+zip routing.</p></body></html>"
    )
    container_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
        '<rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>'
        "</rootfiles></container>"
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        zi = zipfile.ZipInfo("mimetype")
        zi.compress_type = zipfile.ZIP_STORED
        z.writestr(zi, "application/epub+zip")
        z.writestr("META-INF/container.xml", container_xml)
        z.writestr("OEBPS/content.opf", content_opf)
        z.writestr("OEBPS/nav.xhtml", nav_xhtml)
        z.writestr("OEBPS/ch1.xhtml", ch1_xhtml)
    _write("sample.epub", buf.getvalue())


# ── .msg (Outlook Compound File Binary) ───────────────────────────────


def gen_msg() -> None:
    """
    Build a tiny but valid Outlook .msg file (CFB + MAPI).

    Streams emitted:
      __properties_version1.0    — top-level header + property index
      __substg1.0_001A001F       — PR_MESSAGE_CLASS_W ("IPM.Note")
      __substg1.0_0037001F       — PR_SUBJECT_W
      __substg1.0_0E1D001F       — PR_NORMALIZED_SUBJECT_W
      __substg1.0_0C1F001F       — PR_SENDER_EMAIL_ADDRESS_W
      __substg1.0_0042001F       — PR_SENT_REPRESENTING_NAME_W
      __substg1.0_1000001F       — PR_BODY_W
    """
    subject = "owui-cee-proxy routing fixture"
    body = (
        "This is a synthetic .msg fixture used by the proxy's "
        "integration tests. Subject, body and sender are filled "
        "in via MAPI string properties (UTF-16LE)."
    )
    sender_email = "alice@example.com"
    sender_name = "Alice"
    message_class = "IPM.Note"

    string_props = {
        # tag in hex (for the substg name) → string value
        "001A001F": message_class,
        "0037001F": subject,
        "0E1D001F": subject,
        "0C1F001F": sender_email,
        "0042001F": sender_name,
        "1000001F": body,
    }

    streams: dict[str, bytes] = {}
    # __properties_version1.0 top-level header (32 bytes) + 16-byte entry per prop.
    props_stream = bytearray(32)  # all zero per spec for top-level message
    for tag_hex, value in string_props.items():
        encoded = _utf16(value)
        # property entry: tag(4), flags(4), size(4), reserved(4)
        tag = int(tag_hex, 16)
        # Flags: PROPATTR_MANDATORY | PROPATTR_READABLE | PROPATTR_WRITABLE = 0x00000006
        flags = 0x00000006
        # MS-OXMSG: for variable-length strings the size field is the length
        # of the value INCLUDING the trailing UTF-16 NUL (2 bytes).
        size = len(encoded) + 2
        entry = struct.pack("<IIII", tag, flags, size, 0)
        props_stream.extend(entry)
        # The corresponding __substg stream contains the raw UTF-16LE bytes
        # WITHOUT the implicit terminator (terminator is reflected only in
        # the property's size field above).
        streams[f"__substg1.0_{tag_hex}"] = encoded

    streams["__properties_version1.0"] = bytes(props_stream)

    cfb = _build_cfb(streams)
    _write("sample.msg", cfb)


def _utf16(s: str) -> bytes:
    return s.encode("utf-16-le")


def _build_cfb(streams: dict[str, bytes]) -> bytes:
    """
    Minimal v3 (512-byte sectors) Compound-File-Binary writer that
    correctly implements the mini-stream / mini-FAT for streams smaller
    than 4096 bytes — required by olefile and Kreuzberg.

    Layout (logical sectors, header at byte 0..511 is unnumbered):
        sector 0                 : the FAT itself
        sectors 1..D             : directory chain
        sectors D+1..D+M         : mini-FAT chain (M sectors)
        sectors D+M+1..D+M+S     : mini-stream (Root Entry's data, S sectors)

    Each stream lives in mini-sectors (64-byte) inside the mini-stream;
    its directory entry's starting_sector is its mini-sector index
    (0-based) and stream_size is the actual byte length.
    """
    sector_size = 512
    mini_size = 64
    cutoff = 4096
    FAT_FREE = 0xFFFFFFFF
    FAT_END = 0xFFFFFFFE
    FAT_FAT = 0xFFFFFFFD

    header = bytearray(sector_size)
    header[0:8] = bytes.fromhex("D0CF11E0A1B11AE1")            # magic
    struct.pack_into("<H", header, 0x18, 0x003E)               # minor version
    struct.pack_into("<H", header, 0x1A, 0x0003)               # major (v3 = 512B sectors)
    struct.pack_into("<H", header, 0x1C, 0xFFFE)               # byte order LE
    struct.pack_into("<H", header, 0x1E, 9)                    # sector shift 2^9 = 512
    struct.pack_into("<H", header, 0x20, 6)                    # mini sector shift 2^6 = 64

    # Build directory entry helper
    def dir_entry(
        name: str,
        obj_type: int,
        starting_sector: int,
        stream_size: int,
        child_id: int = FAT_FREE,
    ) -> bytes:
        e = bytearray(128)
        nb = name.encode("utf-16-le") + b"\x00\x00"
        if len(nb) > 64:
            raise ValueError(f"name too long: {name}")
        e[0:len(nb)] = nb
        struct.pack_into("<H", e, 0x40, len(nb))
        e[0x42] = obj_type                                      # 2=stream, 5=root
        e[0x43] = 1                                             # node colour: black
        struct.pack_into("<I", e, 0x44, FAT_FREE)              # left sibling
        struct.pack_into("<I", e, 0x48, FAT_FREE)              # right sibling
        struct.pack_into("<I", e, 0x4C, child_id)              # child id
        struct.pack_into("<I", e, 0x74, starting_sector)
        struct.pack_into("<Q", e, 0x78, stream_size)
        return bytes(e)

    # ── Decide which streams live where ───────────────────────────────
    # Streams shorter than `cutoff` go into the mini-stream; rest take
    # full sectors. Per the spec the comparison is strict-less-than.
    mini_streams = {n: d for n, d in streams.items() if len(d) < cutoff}
    full_streams = {n: d for n, d in streams.items() if len(d) >= cutoff}

    # ── Build mini-stream + mini-FAT ──────────────────────────────────
    mini_content = bytearray()
    mini_fat: list[int] = []           # mini-FAT entries, one per mini-sector
    mini_starts: dict[str, int] = {}   # stream name → first mini-sector idx
    next_mini = 0
    for name, data in mini_streams.items():
        n_mini = max(1, (len(data) + mini_size - 1) // mini_size)
        mini_starts[name] = next_mini
        # Pad data to mini-sector boundary then append
        padded = data + b"\x00" * (n_mini * mini_size - len(data))
        mini_content.extend(padded)
        # Mini-FAT chain entries
        for i in range(n_mini):
            this = next_mini + i
            mini_fat.append(FAT_END if i == n_mini - 1 else this + 1)
        next_mini += n_mini
    mini_stream_byte_size = next_mini * mini_size               # exact, unpadded

    # ── Pad mini-FAT to sector multiples ─────────────────────────────
    entries_per_minifat_sector = sector_size // 4
    while len(mini_fat) % entries_per_minifat_sector != 0 or not mini_fat:
        mini_fat.append(FAT_FREE)

    n_minifat_sectors = max(1, len(mini_fat) // entries_per_minifat_sector) if mini_streams else 0
    n_ministream_sectors = (
        (mini_stream_byte_size + sector_size - 1) // sector_size if mini_streams else 0
    )

    # Pad mini_content to whole sectors
    if mini_streams:
        pad = n_ministream_sectors * sector_size - len(mini_content)
        if pad > 0:
            mini_content.extend(b"\x00" * pad)

    # ── Allocate logical sectors ─────────────────────────────────────
    entries_per_sector = sector_size // 128
    n_dir_entries = 1 + len(streams)
    n_dir_sectors = max(1, (n_dir_entries + entries_per_sector - 1) // entries_per_sector)

    dir_sectors = list(range(1, 1 + n_dir_sectors))
    base = 1 + n_dir_sectors
    minifat_sectors = list(range(base, base + n_minifat_sectors))
    base += n_minifat_sectors
    ministream_sectors = list(range(base, base + n_ministream_sectors))
    base += n_ministream_sectors

    # ── Allocate full-sector chains for "large" streams (≥4096) ──────
    full_starts: dict[str, int] = {}
    full_chains: dict[int, list[int]] = {}
    full_data_sectors: dict[int, bytes] = {}
    for name, data in full_streams.items():
        size = len(data)
        n_sectors = max(1, (size + sector_size - 1) // sector_size)
        chain = list(range(base, base + n_sectors))
        base += n_sectors
        full_starts[name] = chain[0]
        full_chains[chain[0]] = chain
        for i, s in enumerate(chain):
            chunk = data[i * sector_size : (i + 1) * sector_size]
            chunk += b"\x00" * (sector_size - len(chunk))
            full_data_sectors[s] = chunk

    next_free = base

    # ── Build the regular FAT ────────────────────────────────────────
    fat = [FAT_FREE] * (sector_size // 4)
    fat[0] = FAT_FAT
    for i, s in enumerate(dir_sectors):
        fat[s] = FAT_END if i == len(dir_sectors) - 1 else dir_sectors[i + 1]
    for i, s in enumerate(minifat_sectors):
        fat[s] = FAT_END if i == len(minifat_sectors) - 1 else minifat_sectors[i + 1]
    for i, s in enumerate(ministream_sectors):
        fat[s] = FAT_END if i == len(ministream_sectors) - 1 else ministream_sectors[i + 1]
    for chain in full_chains.values():
        for i, s in enumerate(chain):
            fat[s] = FAT_END if i == len(chain) - 1 else chain[i + 1]
    fat_sector = b"".join(struct.pack("<I", v) for v in fat)

    # ── Header references ────────────────────────────────────────────
    struct.pack_into("<I", header, 0x28, 0)                    # # dir sectors (v3 = 0)
    struct.pack_into("<I", header, 0x2C, 1)                    # # FAT sectors
    struct.pack_into("<I", header, 0x30, dir_sectors[0])
    struct.pack_into("<I", header, 0x34, 0)
    struct.pack_into("<I", header, 0x38, cutoff)
    struct.pack_into(
        "<I", header, 0x3C,
        minifat_sectors[0] if minifat_sectors else FAT_END,
    )
    struct.pack_into("<I", header, 0x40, n_minifat_sectors)
    struct.pack_into("<I", header, 0x44, FAT_END)              # first DIFAT sector
    struct.pack_into("<I", header, 0x48, 0)                    # # DIFAT sectors
    struct.pack_into("<I", header, 0x4C, 0)                    # DIFAT[0]
    for i in range(1, 109):
        struct.pack_into("<I", header, 0x4C + i * 4, FAT_FREE)

    # Directory sector (sector 1): root + each stream entry.
    # Stream entries MUST be sorted per [MS-CFB] §2.6.4: by name length
    # first, then by uppercase UTF-16LE codepoint comparison. A reader
    # like Kreuzberg that walks the red-black tree will reject the file
    # otherwise.
    def cfb_key(name: str) -> tuple[int, str]:
        return (len(name), name.upper())

    sorted_names = sorted(streams.keys(), key=cfb_key)
    # Map name → directory id (root = 0, then sorted streams).
    dir_id = {name: i + 1 for i, name in enumerate(sorted_names)}

    # Build a degenerate-but-valid red-black tree by inserting entries
    # in sorted order along the right spine: every node's "left" is FREE
    # and "right" points to the next dir id in sorted order. This is a
    # right-leaning chain with all-black colour — valid as long as the
    # comparator on each step holds (which it does because we sort).
    dir_entries: list[bytes] = []

    # Root references the first stream in sorted order; Root.starting_sector
    # = first sector of mini-stream; Root.size = exact mini-stream size.
    first_id = dir_id[sorted_names[0]] if sorted_names else FAT_FREE
    root_start = ministream_sectors[0] if ministream_sectors else FAT_END
    dir_entries.append(
        dir_entry("Root Entry", 5, root_start, mini_stream_byte_size, child_id=first_id)
    )

    for i, name in enumerate(sorted_names):
        # Stream's "starting sector" is interpreted as a mini-sector
        # index when the stream's size is below the cutoff, otherwise
        # as a regular full-sector index.
        if name in mini_starts:
            start = mini_starts[name]
        else:
            start = full_starts[name]
        e = bytearray(dir_entry(name, 2, start, len(streams[name])))
        if i < len(sorted_names) - 1:
            struct.pack_into("<I", e, 0x48, dir_id[sorted_names[i + 1]])
        dir_entries.append(bytes(e))

    # Pad to a whole number of directory sectors.
    needed = n_dir_sectors * entries_per_sector
    while len(dir_entries) < needed:
        empty = bytearray(128)
        empty[0x42] = 0  # unallocated
        struct.pack_into("<I", empty, 0x44, FAT_FREE)
        struct.pack_into("<I", empty, 0x48, FAT_FREE)
        struct.pack_into("<I", empty, 0x4C, FAT_FREE)
        dir_entries.append(bytes(empty))
    dir_blob = b"".join(dir_entries)

    # Materialise the mini-FAT into bytes per sector.
    minifat_blob = bytearray()
    for v in mini_fat:
        minifat_blob.extend(struct.pack("<I", v))

    # Assemble: header + FAT + dir sectors + mini-FAT + mini-stream + full streams.
    body_sectors = [b"\x00" * sector_size] * next_free
    body_sectors[0] = fat_sector
    for i, s in enumerate(dir_sectors):
        body_sectors[s] = dir_blob[i * sector_size : (i + 1) * sector_size]
    for i, s in enumerate(minifat_sectors):
        chunk = bytes(minifat_blob[i * sector_size : (i + 1) * sector_size])
        body_sectors[s] = chunk.ljust(sector_size, b"\x00")
    for i, s in enumerate(ministream_sectors):
        body_sectors[s] = bytes(mini_content[i * sector_size : (i + 1) * sector_size])
    for idx, payload in full_data_sectors.items():
        body_sectors[idx] = payload
    return bytes(header) + b"".join(body_sectors)


# ── README ────────────────────────────────────────────────────────────


def gen_readme() -> None:
    md = """# Integration test fixtures

These files are produced by `scripts/gen-fixtures.py`. Each carries
"owui-cee-proxy routing fixture" somewhere visible so failures and
extracted snippets are unambiguous.

| File           | MIME                                                                                | Routes to (default config)            |
|----------------|--------------------------------------------------------------------------------------|----------------------------------------|
| sample.txt     | text/plain                                                                           | docling                                |
| sample.html    | text/html                                                                            | docling                                |
| sample.md      | text/markdown                                                                        | docling                                |
| sample.csv     | text/csv                                                                             | docling                                |
| sample.json    | application/json                                                                     | docling                                |
| sample.xml     | application/xml                                                                      | docling                                |
| sample.pdf     | application/pdf                                                                      | docling                                |
| sample.png     | image/png                                                                            | docling                                |
| sample.jpeg    | image/jpeg                                                                           | docling                                |
| sample.docx    | application/vnd.openxmlformats-officedocument.wordprocessingml.document              | docling                                |
| sample.xlsx    | application/vnd.openxmlformats-officedocument.spreadsheetml.sheet                    | docling                                |
| sample.pptx    | application/vnd.openxmlformats-officedocument.presentationml.presentation            | docling                                |
| sample.rtf     | application/rtf                                                                      | kreuzberg (fallback)                   |
| sample.eml     | message/rfc822                                                                       | kreuzberg                              |
| sample.odt     | application/vnd.oasis.opendocument.text                                              | kreuzberg                              |
| sample.epub    | application/epub+zip                                                                 | kreuzberg                              |
| sample.msg     | application/vnd.ms-outlook                                                           | kreuzberg                              |

`sample.msg` is a hand-built minimal CFB — enough magic bytes plus
three MAPI string properties (subject / body / sender). It is NOT a
real Outlook export and may not round-trip in Outlook itself, but it
satisfies Kreuzberg's `.msg` parser end-to-end.

To regenerate any of these, run `python3 scripts/gen-fixtures.py`
from the repository root.
"""
    (OUT / "README.md").write_text(md)
    print(f"  wrote {(OUT / 'README.md').relative_to(ROOT)}")


# ── driver ────────────────────────────────────────────────────────────


def main() -> int:
    print(f"Generating fixtures in {OUT.relative_to(ROOT)}/")
    OUT.mkdir(parents=True, exist_ok=True)
    gen_txt()
    gen_html()
    gen_md()
    gen_csv()
    gen_json()
    gen_xml()
    gen_pdf()
    gen_rtf()
    gen_eml()
    gen_png_and_jpeg()
    gen_docx()
    gen_xlsx()
    gen_pptx()
    gen_odt()
    gen_epub()
    gen_msg()
    gen_readme()
    return 0


if __name__ == "__main__":
    sys.exit(main())
